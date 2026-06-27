"""Document parsers for various file formats."""

import bz2
import csv
import gzip
import json
import logging
import lzma
import os
import tarfile
import tempfile
import threading
import xml.etree.ElementTree as ET
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    # pypdf is the maintained successor to PyPDF2; fall back for older installs
    from pypdf import PdfReader
except ImportError:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import yaml
except ImportError:
    yaml = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# ODT/ODS are parsed with the stdlib (zipfile + ElementTree on content.xml); no
# odfpy dependency - it is sdist-only and breaks under Debian's patched setuptools

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    import ebooklib
    from ebooklib import epub
except ImportError:
    epub = None

try:
    import py7zr
except ImportError:
    py7zr = None

try:
    import speech_recognition as sr
except ImportError:
    sr = None

try:
    from pydub import AudioSegment
except ImportError:
    AudioSegment = None

try:
    import pytesseract
    from PIL import Image

    # surface the "binary not on PATH" error class even when mocked in tests
    from pytesseract import TesseractNotFoundError
except ImportError:
    pytesseract = None
    Image = None

    class TesseractNotFoundError(Exception):  # fallback so `except` stays valid
        """Raised when the tesseract-ocr binary is unavailable."""


try:
    # opencv decodes video frames for per-frame OCR
    import cv2
except ImportError:
    cv2 = None

try:
    # faster-whisper: default audio/video speech-to-text backend (CTranslate2)
    from faster_whisper import WhisperModel
except ImportError:
    WhisperModel = None

try:
    # openai-whisper: alternative speech-to-text backend (imports as `whisper`)
    import whisper as openai_whisper
except ImportError:
    openai_whisper = None

logger = logging.getLogger(__name__)

# audio transcription backends, in default preference order
AUDIO_BACKENDS = ("faster-whisper", "openai-whisper", "speech_recognition")


class DocumentParser:
    """Parser for various document formats."""

    # plaintext-ish documents read verbatim
    TEXT_EXTENSIONS = {".txt", ".rst"}

    # source code read as plaintext (niche/example corpora)
    CODE_EXTENSIONS = {
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".java",
        ".kt",
        ".kts",
        ".rs",
        ".c",
        ".h",
        ".cpp",
        ".hpp",
        ".cc",
        ".go",
        ".rb",
        ".php",
        ".swift",
        ".dart",
        ".scala",
        ".sh",
        ".bash",
        ".sql",
        ".r",
        ".m",
        ".pl",
        ".lua",
        ".hs",
        ".clj",
        ".ex",
        ".exs",
        ".vue",
        ".toml",
        ".ini",
        ".cfg",
    }

    # audio formats transcribed to text via speech-to-text
    AUDIO_EXTENSIONS = {
        ".wav",
        ".aiff",
        ".aif",
        ".flac",
        ".mp3",
        ".m4a",
        ".aac",
        ".ogg",
        ".oga",
        ".wma",
        ".opus",
    }

    # speech_recognition reads these natively; others need pydub (+ ffmpeg) to convert
    _NATIVE_AUDIO_EXTENSIONS = {".wav", ".aiff", ".aif", ".flac"}

    # raster images run through OCR (text recognition)
    RASTER_IMAGE_EXTENSIONS = {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".gif",
        ".tif",
        ".tiff",
        ".webp",
        ".ppm",
        ".pgm",
    }

    # svg is vector xml; embedded text is extracted from the markup (no OCR needed)
    VECTOR_IMAGE_EXTENSIONS = {".svg"}

    IMAGE_EXTENSIONS = RASTER_IMAGE_EXTENSIONS | VECTOR_IMAGE_EXTENSIONS

    # video formats: audio track transcribed + on-screen text OCR'd per frame
    VIDEO_EXTENSIONS = {
        ".mp4",
        ".avi",
        ".mov",
        ".mkv",
        ".webm",
        ".flv",
        ".wmv",
        ".mpeg",
        ".mpg",
        ".m4v",
    }

    # extensions of files that can be parsed (and extracted out of archives)
    DOCUMENT_EXTENSIONS = (
        {
            ".md",
            ".pdf",
            ".html",
            ".docx",
            ".csv",
            ".json",
            ".yaml",
            ".yml",
            ".xml",
            ".tex",
            ".xlsx",
            ".pptx",
            ".odt",
            ".ods",
            ".rtf",
            ".epub",
            ".ipynb",
        }
        | TEXT_EXTENSIONS
        | CODE_EXTENSIONS
        | AUDIO_EXTENSIONS
        | IMAGE_EXTENSIONS
        | VIDEO_EXTENSIONS
    )

    # archive containers (extracted, not parsed directly); never recursed into each other
    ARCHIVE_EXTENSIONS = {
        ".zip",
        ".tar",
        ".tar.gz",
        ".tar.bz2",
        ".tar.xz",
        ".tgz",
        ".tbz2",
        ".txz",
        ".gz",
        ".bz2",
        ".xz",
        ".7z",
    }

    SUPPORTED_EXTENSIONS = DOCUMENT_EXTENSIONS | ARCHIVE_EXTENSIONS

    def __init__(
        self,
        audio_backend: str = "faster-whisper",
        whisper_model_size: str = "base",
        ocr_languages: str = "eng",
        video_frame_interval: float = 1.0,
        max_workers: Optional[int] = None,
    ):
        """Initialize the document parser.

        Args:
            audio_backend: speech-to-text backend for audio/video; one of
                "faster-whisper" (default), "openai-whisper", "speech_recognition"
                or "auto". An unavailable backend falls back gracefully.
            whisper_model_size: whisper model size for the whisper backends
                ("tiny", "base", "small", "medium", "large-v3", ...)
            ocr_languages: tesseract language code(s) for image/video OCR (e.g. "eng")
            video_frame_interval: seconds between sampled frames for video OCR
            max_workers: thread-pool size for parse_directory (None -> auto)
        """
        self.audio_backend = audio_backend
        self.whisper_model_size = whisper_model_size
        self.ocr_languages = ocr_languages
        self.video_frame_interval = video_frame_interval
        self.max_workers = max_workers
        # lazily-loaded, reused transcription models (guarded by a lock)
        self._fw_model: Any = None
        self._ow_model: Any = None
        self._model_lock = threading.Lock()
        self._check_dependencies()

    def _check_dependencies(self):
        """Check if optional dependencies are available."""
        missing_deps = []

        if PdfReader is None:
            missing_deps.append("pypdf (for PDF support)")
        if Document is None:
            missing_deps.append("python-docx (for DOCX support)")
        if BeautifulSoup is None:
            missing_deps.append("beautifulsoup4 (for HTML/EPUB support)")
        if yaml is None:
            missing_deps.append("pyyaml (for YAML support)")
        if openpyxl is None:
            missing_deps.append("openpyxl (for XLSX support)")
        if Presentation is None:
            missing_deps.append("python-pptx (for PPTX support)")
        if rtf_to_text is None:
            missing_deps.append("striprtf (for RTF support)")
        if epub is None:
            missing_deps.append("EbookLib (for EPUB support)")
        if py7zr is None:
            missing_deps.append("py7zr (for 7z archive support)")
        if pytesseract is None or Image is None:
            missing_deps.append("pytesseract+Pillow (for image OCR)")
        if cv2 is None:
            missing_deps.append("opencv-python (for video frame OCR)")
        if WhisperModel is None and openai_whisper is None and sr is None:
            missing_deps.append(
                "faster-whisper/openai-whisper/SpeechRecognition "
                "(for audio + video transcription)"
            )

        if missing_deps:
            logger.warning(f"Missing optional dependencies: {', '.join(missing_deps)}")

    def parse_directory(
        self, directory_path: str, max_workers: Optional[int] = None
    ) -> List[Dict[str, Any]]:
        """
        Recursively parse all supported documents in a directory.

        Files are parsed concurrently with a thread pool; parsing is largely I/O
        and native-library bound (PDF, OCR, transcription), so threads overlap
        well. The returned list is sorted by path for stable ordering regardless
        of completion order.

        Args:
            directory_path: Path to the directory to scan
            max_workers: Thread-pool size (None -> self.max_workers -> auto)

        Returns:
            List of parsed documents with metadata
        """
        directory_path = Path(directory_path)

        if not directory_path.exists():
            raise FileNotFoundError(f"Directory not found: {directory_path}")

        # Recursively find all supported files
        files = [
            fp
            for fp in directory_path.rglob("*")
            if fp.is_file() and self._resolve_extension(fp) in self.SUPPORTED_EXTENSIONS
        ]

        workers = self._resolve_workers(
            max_workers if max_workers is not None else self.max_workers, len(files)
        )

        documents: List[Dict[str, Any]] = []
        if workers <= 1:
            for file_path in files:
                doc = self._safe_parse_file(file_path, directory_path)
                if doc:
                    documents.append(doc)
        else:
            logger.info(f"Parsing {len(files)} files with {workers} worker threads")
            with ThreadPoolExecutor(max_workers=workers) as executor:
                futures = [
                    executor.submit(self._safe_parse_file, fp, directory_path)
                    for fp in files
                ]
                for future in as_completed(futures):
                    doc = future.result()
                    if doc:
                        documents.append(doc)

        # stable order regardless of thread completion order
        documents.sort(key=lambda d: d["filepath"])
        return documents

    def _safe_parse_file(
        self, file_path: Path, base_path: Path
    ) -> Optional[Dict[str, Any]]:
        """parse_file wrapper that logs and swallows per-file errors."""
        try:
            return self.parse_file(file_path, base_path=base_path)
        except Exception as e:
            logger.error(f"Error parsing {file_path}: {e}")
            return None

    @staticmethod
    def _resolve_workers(max_workers: Optional[int], num_files: int) -> int:
        """Resolve worker count (auto: ~cpu count, capped at 8 and <= file count)."""
        if max_workers is not None:
            return max(1, max_workers)
        if num_files <= 1:
            return 1
        cpu = os.cpu_count() or 1
        # modest default; parsing is I/O/native bound, cap to avoid oversubscription
        return max(1, min(8, cpu, num_files))

    def _resolve_extension(self, file_path: Path) -> str:
        """Resolve an extension, accounting for compound tar suffixes."""
        name = file_path.name.lower()
        if name.endswith((".tar.gz", ".tgz")):
            return ".tar.gz"
        if name.endswith((".tar.bz2", ".tbz2")):
            return ".tar.bz2"
        if name.endswith((".tar.xz", ".txz")):
            return ".tar.xz"
        return file_path.suffix.lower()

    def parse_file(
        self, file_path: Path, base_path: Optional[Path] = None
    ) -> Optional[Dict[str, Any]]:
        """
        Parse a single file based on its extension.

        Args:
            file_path: Path to the file to parse
            base_path: Base directory path for calculating relative labels

        Returns:
            Parsed document with metadata or None if parsing failed
        """
        extension = self._resolve_extension(file_path)

        try:
            if extension == ".zip":
                content = self._parse_zip(file_path)
            elif extension in (".tar", ".tar.gz", ".tar.bz2", ".tar.xz"):
                content = self._parse_tar(file_path)
            elif extension == ".7z":
                content = self._parse_7z(file_path)
            elif extension in (".gz", ".bz2", ".xz"):
                content = self._parse_compressed(file_path)
            else:
                content = self._content_for_extension(file_path, extension)

            if content is None:
                logger.warning(f"Unsupported file type: {extension}")
                return None

            if content:
                # Calculate directory-based labels if base_path is provided
                filename_stem = file_path.stem  # filename without extension

                if base_path:
                    try:
                        relative_path = file_path.relative_to(base_path)
                        # Combine parent directory and filename
                        if (
                            relative_path.parent.name
                            and relative_path.parent.name != "."
                        ):
                            dir_name = relative_path.parent.name
                            label = f"{dir_name}_{filename_stem}"
                            category_path = str(relative_path.parent).replace("\\", "/")
                        else:
                            label = f"root_{filename_stem}"
                            category_path = ""
                    except ValueError:
                        # file_path is not relative to base_path
                        dir_name = (
                            file_path.parent.name
                            if file_path.parent.name
                            else "unknown"
                        )
                        label = f"{dir_name}_{filename_stem}"
                        category_path = ""
                else:
                    dir_name = (
                        file_path.parent.name if file_path.parent.name else "unknown"
                    )
                    label = f"{dir_name}_{filename_stem}"
                    category_path = ""

                return {
                    "content": content,
                    "filename": file_path.name,
                    "filepath": str(file_path),
                    "extension": extension,
                    "size": file_path.stat().st_size,
                    "label": label,
                    "category_path": category_path,
                }

        except Exception as e:
            logger.error(f"Error parsing {file_path}: {e}")
            return None

    def _content_for_extension(self, file_path: Path, extension: str) -> Optional[str]:
        """Dispatch a single (already-resolved) document extension to its parser."""
        if extension == ".md":
            return self._parse_markdown(file_path)
        if (
            extension == ""
            or extension in self.TEXT_EXTENSIONS
            or extension in self.CODE_EXTENSIONS
        ):
            return self._parse_text(file_path)
        if extension == ".tex":
            return self._parse_latex(file_path)
        if extension == ".pdf":
            return self._parse_pdf(file_path)
        if extension == ".html":
            return self._parse_html(file_path)
        if extension == ".docx":
            return self._parse_docx(file_path)
        if extension == ".csv":
            return self._parse_csv(file_path)
        if extension == ".json":
            return self._parse_json(file_path)
        if extension == ".ipynb":
            return self._parse_ipynb(file_path)
        if extension in (".yaml", ".yml"):
            return self._parse_yaml(file_path)
        if extension == ".xml":
            return self._parse_xml(file_path)
        if extension == ".xlsx":
            return self._parse_xlsx(file_path)
        if extension == ".pptx":
            return self._parse_pptx(file_path)
        if extension == ".odt":
            return self._parse_odt(file_path)
        if extension == ".ods":
            return self._parse_ods(file_path)
        if extension == ".rtf":
            return self._parse_rtf(file_path)
        if extension == ".epub":
            return self._parse_epub(file_path)
        if extension in self.AUDIO_EXTENSIONS:
            return self._transcribe_media(file_path)
        if extension in self.VECTOR_IMAGE_EXTENSIONS:
            return self._parse_svg(file_path)
        if extension in self.RASTER_IMAGE_EXTENSIONS:
            return self._parse_image(file_path)
        if extension in self.VIDEO_EXTENSIONS:
            return self._parse_video(file_path)
        return None

    def _parse_markdown(self, file_path: Path) -> str:
        """Parse Markdown file."""
        return self._parse_text(file_path)

    def _parse_text(self, file_path: Path) -> str:
        """Parse text/code file; skip gracefully if it is not utf-8 decodable."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            logger.warning(f"Skipping non-text (binary?) file: {file_path}")
            return ""

    def _parse_pdf(self, file_path: Path) -> str:
        """Parse PDF file."""
        if PdfReader is None:
            logger.error("pypdf/PyPDF2 not installed. Cannot parse PDF files.")
            return ""

        text = ""
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _parse_html(self, file_path: Path) -> str:
        """Parse HTML file."""
        if BeautifulSoup is None:
            logger.error("BeautifulSoup not installed. Cannot parse HTML files.")
            return ""

        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            return soup.get_text()

    def _parse_docx(self, file_path: Path) -> str:
        """Parse Word document."""
        if Document is None:
            logger.error("python-docx not installed. Cannot parse DOCX files.")
            return ""

        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    def _parse_csv(self, file_path: Path) -> str:
        """Parse CSV file."""
        text = ""
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                text += ", ".join(row) + "\n"
        return text

    def _parse_json(self, file_path: Path) -> str:
        """Parse JSON file."""
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
            return json.dumps(data, indent=2)

    def _parse_ipynb(self, file_path: Path) -> str:
        """Parse a Jupyter notebook; join markdown and code cells."""
        with open(file_path, "r", encoding="utf-8") as f:
            notebook = json.load(f)

        parts = []
        for cell in notebook.get("cells", []):
            source = cell.get("source", "")
            if isinstance(source, list):
                source = "".join(source)
            if not source.strip():
                continue
            if cell.get("cell_type") == "code":
                parts.append(f"```\n{source}\n```")
            else:
                parts.append(source)
        return "\n\n".join(parts)

    def _parse_yaml(self, file_path: Path) -> str:
        """Parse YAML file."""
        if yaml is None:
            logger.error("PyYAML not installed. Cannot parse YAML files.")
            return ""

        with open(file_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
            return yaml.dump(data, default_flow_style=False)

    def _parse_xml(self, file_path: Path) -> str:
        """Parse XML file."""
        with open(file_path, "r", encoding="utf-8") as f:
            tree = ET.parse(f)
            root = tree.getroot()
            return ET.tostring(root, encoding="unicode")

    def _parse_xlsx(self, file_path: Path) -> str:
        """Parse Excel (XLSX) file."""
        if openpyxl is None:
            logger.error("openpyxl not installed. Cannot parse XLSX files.")
            return ""

        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"Sheet: {sheet_name}")

                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        if cell.value is not None:
                            row_data.append(str(cell.value))
                    if row_data:  # Only add non-empty rows
                        text_content.append("\t".join(row_data))

                text_content.append("")  # Empty line between sheets

            return "\n".join(text_content)

        except Exception as e:
            logger.error(f"Error parsing XLSX file {file_path}: {e}")
            return ""

    def _parse_pptx(self, file_path: Path) -> str:
        """Parse a PowerPoint (PPTX) file; slide text plus speaker notes."""
        if Presentation is None:
            logger.error("python-pptx not installed. Cannot parse PPTX files.")
            return ""

        try:
            prs = Presentation(str(file_path))
            parts = []
            for index, slide in enumerate(prs.slides, start=1):
                parts.append(f"Slide {index}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs)
                            if line.strip():
                                parts.append(line)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        parts.append(f"[notes] {notes}")
                parts.append("")
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}")
            return ""

    def _parse_odt(self, file_path: Path) -> str:
        """Parse an OpenDocument text (ODT) file (stdlib; reads content.xml)."""
        return self._parse_opendocument(file_path, spreadsheet=False)

    def _parse_ods(self, file_path: Path) -> str:
        """Parse an OpenDocument spreadsheet (ODS) file (stdlib; reads content.xml)."""
        return self._parse_opendocument(file_path, spreadsheet=True)

    def _parse_opendocument(self, file_path: Path, spreadsheet: bool) -> str:
        """Extract text from an ODT/ODS by reading content.xml (no odfpy needed).

        OpenDocument files are zips; content.xml holds the body. Pull text from
        text:p / text:h elements (and table cells for spreadsheets) - enough for a
        training corpus without the sdist-only, unmaintained odfpy dependency.
        """
        text_ns = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"
        table_ns = "{urn:oasis:names:tc:opendocument:xmlns:table:1.0}"
        try:
            with zipfile.ZipFile(file_path) as zf:
                data = zf.read("content.xml")
            root = ET.fromstring(data)
        except (KeyError, zipfile.BadZipFile, ET.ParseError, OSError) as e:
            logger.error(f"Error reading OpenDocument file {file_path}: {e}")
            return ""

        if not spreadsheet:
            # text doc: collect paragraphs and headings in document order
            parts = []
            for elem in root.iter():
                if elem.tag in (text_ns + "p", text_ns + "h"):
                    text = "".join(elem.itertext())
                    if text.strip():
                        parts.append(text)
            return "\n".join(parts)

        # spreadsheet: walk tables -> rows -> cells
        parts = []
        for table in root.iter(table_ns + "table"):
            name = table.get(table_ns + "name")
            if name:
                parts.append(f"Sheet: {name}")
            for row in table.iter(table_ns + "table-row"):
                cells = []
                for cell in row.iter(table_ns + "table-cell"):
                    cell_text = "".join(cell.itertext()).strip()
                    if cell_text:
                        cells.append(cell_text)
                if cells:
                    parts.append("\t".join(cells))
            parts.append("")
        return "\n".join(parts)

    def _parse_rtf(self, file_path: Path) -> str:
        """Parse a Rich Text Format (RTF) file."""
        if rtf_to_text is None:
            logger.error("striprtf not installed. Cannot parse RTF files.")
            return ""

        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return rtf_to_text(f.read())
        except Exception as e:
            logger.error(f"Error parsing RTF file {file_path}: {e}")
            return ""

    def _parse_epub(self, file_path: Path) -> str:
        """Parse an EPUB e-book; extract chapter text."""
        if epub is None or BeautifulSoup is None:
            logger.error("EbookLib/beautifulsoup4 not installed. Cannot parse EPUB.")
            return ""

        try:
            book = epub.read_epub(str(file_path))
            parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text()
                if text.strip():
                    parts.append(text)
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"Error parsing EPUB file {file_path}: {e}")
            return ""

    # ---- audio / video transcription (speech-to-text) ----

    def _resolve_audio_backend(self) -> Optional[str]:
        """Resolve the effective transcription backend, honoring availability.

        Honors the configured preference first (faster-whisper by default), then
        falls back through the remaining backends. "auto" just takes the first
        available in preference order.
        """
        available = {
            "faster-whisper": WhisperModel is not None,
            "openai-whisper": openai_whisper is not None,
            "speech_recognition": sr is not None,
        }
        pref = (self.audio_backend or "faster-whisper").lower()
        if pref == "auto":
            order = list(AUDIO_BACKENDS)
        else:
            order = [pref] + [b for b in AUDIO_BACKENDS if b != pref]

        for backend in order:
            if available.get(backend):
                if pref not in ("auto", backend):
                    logger.warning(
                        f"Audio backend '{pref}' unavailable; using '{backend}'"
                    )
                return backend

        logger.error(
            "No transcription backend available. Install faster-whisper, "
            "openai-whisper, or SpeechRecognition."
        )
        return None

    def _whisper_device_compute(self):
        """Pick (device, compute_type) for the whisper backends."""
        try:
            import torch

            if torch.cuda.is_available():
                return "cuda", "float16"
        except Exception:
            pass
        # int8 keeps CPU transcription fast and light (matches faster-whisper docs)
        return "cpu", "int8"

    def _transcribe_media(self, file_path: Path) -> str:
        """Transcribe speech from an audio or video file via the chosen backend."""
        backend = self._resolve_audio_backend()
        if backend == "faster-whisper":
            return self._transcribe_faster_whisper(file_path)
        if backend == "openai-whisper":
            return self._transcribe_openai_whisper(file_path)
        if backend == "speech_recognition":
            return self._transcribe_speech_recognition(file_path)
        return ""

    def _transcribe_faster_whisper(self, file_path: Path) -> str:
        """Transcribe via faster-whisper (CTranslate2); reads media containers."""
        if WhisperModel is None:
            return ""
        try:
            if self._fw_model is None:
                # double-checked lock so the model loads once across threads
                with self._model_lock:
                    if self._fw_model is None:
                        device, compute_type = self._whisper_device_compute()
                        logger.info(
                            f"Loading faster-whisper '{self.whisper_model_size}' "
                            f"on {device} ({compute_type})"
                        )
                        self._fw_model = WhisperModel(
                            self.whisper_model_size,
                            device=device,
                            compute_type=compute_type,
                        )
            # CTranslate2 transcription is thread-safe; no lock around inference
            segments, _info = self._fw_model.transcribe(str(file_path))
            return " ".join(segment.text.strip() for segment in segments).strip()
        except Exception as e:
            logger.error(f"faster-whisper failed on {file_path}: {e}")
            return ""

    def _transcribe_openai_whisper(self, file_path: Path) -> str:
        """Transcribe via openai-whisper (torch); reads media containers."""
        if openai_whisper is None:
            return ""
        try:
            # torch models are not thread-safe; serialize load + inference
            with self._model_lock:
                if self._ow_model is None:
                    logger.info(f"Loading openai-whisper '{self.whisper_model_size}'")
                    self._ow_model = openai_whisper.load_model(self.whisper_model_size)
                result = self._ow_model.transcribe(str(file_path))
            return (result.get("text") or "").strip()
        except Exception as e:
            logger.error(f"openai-whisper failed on {file_path}: {e}")
            return ""

    def _transcribe_speech_recognition(self, file_path: Path) -> str:
        """Legacy fallback: transcribe via SpeechRecognition (Google Web Speech).

        WAV/AIFF/FLAC are read natively; other formats (mp3, m4a, video, ...) are
        converted to wav first with pydub (which needs the ffmpeg binary).
        """
        if sr is None:
            logger.error(
                "SpeechRecognition not installed. Cannot transcribe audio files."
            )
            return ""

        extension = file_path.suffix.lower()
        tmp_wav: Optional[Path] = None

        try:
            if extension in self._NATIVE_AUDIO_EXTENSIONS:
                audio_path = file_path
            else:
                if AudioSegment is None:
                    logger.error(
                        f"pydub (+ ffmpeg) needed to transcribe '{extension}' media."
                    )
                    return ""
                with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as handle:
                    tmp_wav = Path(handle.name)
                AudioSegment.from_file(str(file_path)).export(
                    str(tmp_wav), format="wav"
                )
                audio_path = tmp_wav

            recognizer = sr.Recognizer()
            with sr.AudioFile(str(audio_path)) as source:
                audio = recognizer.record(source)

            try:
                return recognizer.recognize_google(audio)
            except sr.UnknownValueError:
                logger.warning(f"Speech not understood in audio: {file_path}")
                return ""
            except sr.RequestError as e:
                logger.error(f"Speech-to-text request failed for {file_path}: {e}")
                return ""

        except Exception as e:
            logger.error(f"Error transcribing media {file_path}: {e}")
            return ""
        finally:
            if tmp_wav is not None:
                tmp_wav.unlink(missing_ok=True)

    # ---- image text recognition (OCR) ----

    def _parse_image(self, file_path: Path) -> str:
        """Recognize text in a raster image via OCR (tesseract)."""
        if pytesseract is None or Image is None:
            logger.error("pytesseract+Pillow not installed. Cannot OCR images.")
            return ""
        try:
            with Image.open(file_path) as img:
                return pytesseract.image_to_string(img, lang=self.ocr_languages).strip()
        except TesseractNotFoundError:
            logger.error(
                "tesseract-ocr binary not found on PATH. Install it to OCR images."
            )
            return ""
        except Exception as e:
            logger.error(f"Error running OCR on image {file_path}: {e}")
            return ""

    def _parse_svg(self, file_path: Path) -> str:
        """Extract embedded text from an SVG (vector xml; no OCR needed)."""
        try:
            tree = ET.parse(file_path)
            parts = []
            for elem in tree.getroot().iter():
                tag = elem.tag.rsplit("}", 1)[-1]  # strip xml namespace
                if tag in ("text", "tspan", "title", "desc"):
                    if elem.text and elem.text.strip():
                        parts.append(elem.text.strip())
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"Error parsing SVG {file_path}: {e}")
            return ""

    # ---- video (audio transcript + per-frame OCR) ----

    def _parse_video(self, file_path: Path) -> str:
        """Parse a video: transcribe the audio track and OCR on-screen text.

        The audio track is transcribed via the whisper backend (which reads the
        video container directly); frames are sampled and OCR'd, deduped by
        identical text so a slide shown for many frames is captured once.
        """
        parts = []

        transcript = self._transcribe_media(file_path)
        if transcript.strip():
            parts.append("=== Speech transcript ===")
            parts.append(transcript.strip())

        on_screen = self._ocr_video_frames(file_path)
        if on_screen.strip():
            parts.append("=== On-screen text ===")
            parts.append(on_screen.strip())

        return "\n\n".join(parts)

    def _ocr_video_frames(self, file_path: Path) -> str:
        """Sample video frames and OCR them, deduping identical on-screen text."""
        if cv2 is None or pytesseract is None:
            logger.error("opencv-python+pytesseract needed for video frame OCR.")
            return ""

        capture = cv2.VideoCapture(str(file_path))
        if not capture.isOpened():
            logger.error(f"Could not open video for OCR: {file_path}")
            return ""

        try:
            fps = capture.get(cv2.CAP_PROP_FPS) or 0.0
            # sample once per `video_frame_interval` seconds (fallback: every frame)
            interval = int(round(fps * self.video_frame_interval)) if fps > 0 else 1
            interval = max(1, interval)

            parts: List[str] = []
            seen = set()  # dedupe identical on-screen text across frames
            frame_idx = 0
            while True:
                ret, frame = capture.read()
                if not ret:
                    break
                if frame_idx % interval == 0:
                    gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                    try:
                        text = pytesseract.image_to_string(
                            gray, lang=self.ocr_languages
                        ).strip()
                    except TesseractNotFoundError:
                        logger.error("tesseract-ocr binary not found on PATH.")
                        break
                    normalized = " ".join(text.split())
                    if normalized and normalized not in seen:
                        seen.add(normalized)
                        parts.append(text)
                frame_idx += 1
            return "\n".join(parts)
        finally:
            capture.release()

    def _parse_latex(self, file_path: Path) -> str:
        """Parse LaTeX file."""
        return self._parse_text(file_path)

    def _parse_zip(self, file_path: Path) -> str:
        """Parse ZIP archive by extracting and parsing supported documents."""
        content_parts = []

        try:
            with zipfile.ZipFile(file_path, "r") as zip_ref:
                # Get list of files in the archive
                file_list = zip_ref.namelist()
                content_parts.append(f"=== ZIP Archive: {file_path.name} ===")
                content_parts.append(f"Contains {len(file_list)} files:")

                # Create a temporary directory to extract files
                with tempfile.TemporaryDirectory() as temp_dir:
                    for file_info in zip_ref.infolist():
                        # Skip directories and hidden files
                        if file_info.is_dir() or file_info.filename.startswith("."):
                            continue

                        file_name = Path(file_info.filename)
                        extension = file_name.suffix.lower()

                        # Only process files with supported document extensions
                        if extension in self.DOCUMENT_EXTENSIONS:
                            try:
                                # Extract the file to temp directory
                                extracted_path = zip_ref.extract(file_info, temp_dir)
                                extracted_file = Path(extracted_path)

                                # Parse the extracted file
                                parsed_doc = self._parse_extracted_file(
                                    extracted_file, file_info.filename
                                )
                                if parsed_doc:
                                    content_parts.append(
                                        f"\n--- File: {file_info.filename} ---"
                                    )
                                    content_parts.append(parsed_doc)

                            except Exception as e:
                                logger.warning(
                                    f"Could not parse {file_info.filename} from ZIP: {e}"
                                )
                                content_parts.append(
                                    f"\n--- File: {file_info.filename} (parsing failed) ---"
                                )
                        else:
                            content_parts.append(
                                f"  • {file_info.filename} (unsupported format)"
                            )

        except Exception as e:
            logger.error(f"Error reading ZIP file {file_path}: {e}")
            return f"Error reading ZIP file: {e}"

        return "\n".join(content_parts)

    def _parse_tar(self, file_path: Path) -> str:
        """Parse TAR archive by extracting and parsing supported documents."""
        content_parts = []

        try:
            # Determine the compression mode
            if file_path.name.lower().endswith((".tar.gz", ".tgz")):
                mode = "r:gz"
            elif file_path.name.lower().endswith((".tar.bz2", ".tbz2")):
                mode = "r:bz2"
            elif file_path.name.lower().endswith((".tar.xz", ".txz")):
                mode = "r:xz"
            else:
                mode = "r"

            with tarfile.open(file_path, mode) as tar_ref:
                # Get list of files in the archive
                members = tar_ref.getmembers()
                file_members = [m for m in members if m.isfile()]

                content_parts.append(f"=== TAR Archive: {file_path.name} ===")
                content_parts.append(f"Contains {len(file_members)} files:")

                # Create a temporary directory to extract files
                with tempfile.TemporaryDirectory() as temp_dir:
                    temp_path = Path(temp_dir)

                    for member in file_members:
                        # Skip hidden files
                        if Path(member.name).name.startswith("."):
                            continue

                        file_name = Path(member.name)
                        extension = file_name.suffix.lower()

                        # Only process files with supported document extensions
                        if extension in self.DOCUMENT_EXTENSIONS:
                            try:
                                # Extract the file to temp directory
                                tar_ref.extract(member, temp_dir)
                                extracted_file = temp_path / member.name

                                # Parse the extracted file
                                parsed_doc = self._parse_extracted_file(
                                    extracted_file, member.name
                                )
                                if parsed_doc:
                                    content_parts.append(
                                        f"\n--- File: {member.name} ---"
                                    )
                                    content_parts.append(parsed_doc)

                            except Exception as e:
                                logger.warning(
                                    f"Could not parse {member.name} from TAR: {e}"
                                )
                                content_parts.append(
                                    f"\n--- File: {member.name} (parsing failed) ---"
                                )
                        else:
                            content_parts.append(
                                f"  • {member.name} (unsupported format)"
                            )

        except Exception as e:
            logger.error(f"Error reading TAR file {file_path}: {e}")
            return f"Error reading TAR file: {e}"

        return "\n".join(content_parts)

    def _parse_7z(self, file_path: Path) -> str:
        """Parse a 7z archive by extracting and parsing supported documents."""
        if py7zr is None:
            logger.error("py7zr not installed. Cannot parse 7z archives.")
            return ""

        content_parts = [f"=== 7z Archive: {file_path.name} ==="]

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                with py7zr.SevenZipFile(file_path, "r") as archive:
                    archive.extractall(path=temp_dir)

                temp_path = Path(temp_dir)
                for extracted_file in sorted(temp_path.rglob("*")):
                    if not extracted_file.is_file():
                        continue
                    if extracted_file.name.startswith("."):
                        continue
                    if extracted_file.suffix.lower() not in self.DOCUMENT_EXTENSIONS:
                        content_parts.append(
                            f"  • {extracted_file.name} (unsupported format)"
                        )
                        continue

                    rel = extracted_file.relative_to(temp_path)
                    parsed_doc = self._parse_extracted_file(extracted_file, str(rel))
                    if parsed_doc:
                        content_parts.append(f"\n--- File: {rel} ---")
                        content_parts.append(parsed_doc)

        except Exception as e:
            logger.error(f"Error reading 7z file {file_path}: {e}")
            return f"Error reading 7z file: {e}"

        return "\n".join(content_parts)

    def _parse_compressed(self, file_path: Path) -> str:
        """Decompress a single-file .gz/.bz2/.xz and parse its inner document."""
        name = file_path.name.lower()
        if name.endswith(".gz"):
            opener = gzip.open
        elif name.endswith(".bz2"):
            opener = bz2.open
        elif name.endswith(".xz"):
            opener = lzma.open
        else:
            return ""

        # inner filename is the name minus the compression suffix
        inner_name = file_path.name[: -(len(file_path.suffix))]
        inner_ext = Path(inner_name).suffix.lower()

        try:
            with opener(file_path, "rb") as f:
                data = f.read()

            with tempfile.TemporaryDirectory() as temp_dir:
                inner_path = Path(temp_dir) / (inner_name or "decompressed.txt")
                inner_path.write_bytes(data)
                # treat unknown inner extensions as plaintext
                if inner_ext not in self.DOCUMENT_EXTENSIONS:
                    return self._parse_text(inner_path)
                parsed = self._parse_extracted_file(inner_path, inner_name)
                return parsed or ""
        except Exception as e:
            logger.error(f"Error reading compressed file {file_path}: {e}")
            return f"Error reading compressed file: {e}"

    def _parse_extracted_file(
        self, file_path: Path, original_name: str
    ) -> Optional[str]:
        """Parse an extracted file from an archive (no nested archive recursion)."""
        extension = file_path.suffix.lower()

        try:
            return self._content_for_extension(file_path, extension)
        except Exception as e:
            logger.warning(f"Error parsing extracted file {original_name}: {e}")
            return None
