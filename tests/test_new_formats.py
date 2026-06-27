"""Tests for the document formats added in v1.0.0."""

import gzip
import json

import pytest

from doc2lora.parsers import DocumentParser


@pytest.fixture
def parser():
    return DocumentParser()


# ---- formats with no optional dependency ----


def test_parse_rst(parser, temp_dir):
    """reStructuredText is read as plaintext."""
    content = "Title\n=====\n\nSome rst body text."
    f = temp_dir / "doc.rst"
    f.write_text(content)

    result = parser.parse_file(f)

    assert result is not None
    assert result["content"] == content
    assert result["extension"] == ".rst"


@pytest.mark.parametrize("ext", [".py", ".rs", ".kt", ".cpp", ".go", ".dart"])
def test_parse_code_as_plaintext(parser, temp_dir, ext):
    """Source code files are read verbatim as plaintext."""
    content = "fn main() { print('hi') }\n# sample code"
    f = temp_dir / f"sample{ext}"
    f.write_text(content)

    result = parser.parse_file(f)

    assert result is not None
    assert result["content"] == content
    assert result["extension"] == ext


def test_parse_ipynb(parser, temp_dir):
    """Jupyter notebooks join markdown and code cells."""
    notebook = {
        "cells": [
            {"cell_type": "markdown", "source": ["# Heading\n", "intro text"]},
            {"cell_type": "code", "source": "print('hello')"},
            {"cell_type": "code", "source": "   "},  # blank skipped
        ],
        "metadata": {},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    f = temp_dir / "nb.ipynb"
    f.write_text(json.dumps(notebook))

    result = parser.parse_file(f)

    assert result is not None
    assert "# Heading" in result["content"]
    assert "print('hello')" in result["content"]
    assert "```" in result["content"]  # code cells get fenced


def test_binary_file_with_code_extension_is_skipped(parser, temp_dir):
    """A binary file mislabeled with a code extension degrades gracefully."""
    f = temp_dir / "weird.py"
    f.write_bytes(b"\x80\x81\x82\xff\xfe")

    result = parser.parse_file(f)

    assert result is None  # empty content -> skipped, not a crash


def test_single_file_gzip(parser, temp_dir):
    """A single-file .gz is decompressed and its inner document parsed."""
    f = temp_dir / "notes.txt.gz"
    with gzip.open(f, "wb") as fh:
        fh.write(b"decompressed body")

    result = parser.parse_file(f)

    assert result is not None
    assert "decompressed body" in result["content"]
    assert result["extension"] == ".gz"


# ---- formats gated on optional dependencies ----


def test_parse_pptx(parser, temp_dir):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "presentation body text"
    f = temp_dir / "deck.pptx"
    prs.save(str(f))

    result = parser.parse_file(f)

    assert result is not None
    assert "presentation body text" in result["content"]


def _write_odf(path, content_xml):
    """Write a minimal OpenDocument file (a zip whose content.xml holds the body)."""
    import zipfile

    with zipfile.ZipFile(path, "w") as z:
        z.writestr("content.xml", content_xml)


def test_parse_odt(parser, temp_dir):
    """ODT text is extracted from content.xml with the stdlib (no odfpy)."""
    content_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<office:document-content"
        ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
        ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
        "<office:body><office:text>"
        "<text:h>Heading One</text:h>"
        "<text:p>opendocument body text</text:p>"
        "</office:text></office:body></office:document-content>"
    )
    f = temp_dir / "doc.odt"
    _write_odf(f, content_xml)

    result = parser.parse_file(f)

    assert result is not None
    assert "opendocument body text" in result["content"]
    assert "Heading One" in result["content"]
    assert result["extension"] == ".odt"


def test_parse_ods(parser, temp_dir):
    """ODS cell text is extracted from content.xml with the stdlib (no odfpy)."""
    content_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<office:document-content"
        ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
        ' xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0"'
        ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
        "<office:body><office:spreadsheet>"
        '<table:table table:name="Sheet1">'
        "<table:table-row>"
        "<table:table-cell><text:p>cellvalue</text:p></table:table-cell>"
        "<table:table-cell><text:p>42</text:p></table:table-cell>"
        "</table:table-row>"
        "</table:table>"
        "</office:spreadsheet></office:body></office:document-content>"
    )
    f = temp_dir / "sheet.ods"
    _write_odf(f, content_xml)

    result = parser.parse_file(f)

    assert result is not None
    assert "cellvalue" in result["content"]
    assert "Sheet1" in result["content"]
    assert result["extension"] == ".ods"


def test_parse_rtf(parser, temp_dir):
    pytest.importorskip("striprtf")
    rtf = r"{\rtf1\ansi rich text body\par}"
    f = temp_dir / "doc.rtf"
    f.write_text(rtf)

    result = parser.parse_file(f)

    assert result is not None
    assert "rich text body" in result["content"]


def test_parse_epub(parser, temp_dir):
    pytest.importorskip("ebooklib")
    pytest.importorskip("bs4")
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("id1")
    book.set_title("Test Book")
    book.set_language("en")
    chapter = epub.EpubHtml(title="C1", file_name="c1.xhtml", lang="en")
    chapter.content = "<html><body><p>epub chapter body</p></body></html>"
    book.add_item(chapter)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chapter]
    f = temp_dir / "book.epub"
    epub.write_epub(str(f), book)

    result = parser.parse_file(f)

    assert result is not None
    assert "epub chapter body" in result["content"]


def test_parse_audio_wav(temp_dir, monkeypatch):
    """A WAV file is transcribed via the SpeechRecognition backend (mocked)."""
    sr = pytest.importorskip("speech_recognition")
    import wave

    # pin the legacy backend so the test is deterministic regardless of which
    # whisper backends are installed (faster-whisper is the default)
    parser = DocumentParser(audio_backend="speech_recognition")

    wav = temp_dir / "clip.wav"
    with wave.open(str(wav), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(16000)
        w.writeframes(b"\x00\x00" * 16000)  # 1s of silence

    # avoid the network call to Google's web speech API
    monkeypatch.setattr(
        sr.Recognizer,
        "recognize_google",
        lambda self, audio, **kwargs: "transcribed audio text",
    )

    result = parser.parse_file(wav)

    assert result is not None
    assert "transcribed audio text" in result["content"]
    assert result["extension"] == ".wav"


def test_audio_extensions_registered(parser):
    """Common audio formats are recognized as supported documents."""
    for ext in (".mp3", ".m4a", ".flac", ".aac", ".wav"):
        assert ext in parser.DOCUMENT_EXTENSIONS
        assert ext in parser.AUDIO_EXTENSIONS


def test_parse_7z(parser, temp_dir):
    py7zr = pytest.importorskip("py7zr")
    inner = temp_dir / "inner.txt"
    inner.write_text("seven zip inner content")
    f = temp_dir / "archive.7z"
    with py7zr.SevenZipFile(f, "w") as archive:
        archive.write(str(inner), "inner.txt")

    result = parser.parse_file(f)

    assert result is not None
    assert "seven zip inner content" in result["content"]
    assert result["extension"] == ".7z"
